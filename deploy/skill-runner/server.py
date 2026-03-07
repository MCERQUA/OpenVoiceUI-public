"""
skill-runner/server.py — Shared JamBot skill execution service

Runs on jambot-shared Docker network at http://skill-runner:8900
NOT exposed to the host or the internet — internal only.

Endpoints:
  GET  /health          — liveness check
  POST /extract         — document text extraction (PDF, DOCX, XLSX, PPTX)
  POST /analyze-csv     — pandas CSV summary + stats
"""

import logging
import os
import re
import tempfile
from pathlib import Path

from flask import Flask, jsonify, request

logging.basicConfig(level=logging.INFO, format='%(asctime)s %(levelname)s %(message)s')
logger = logging.getLogger(__name__)

app = Flask(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

MAX_FILE_BYTES = 25 * 1024 * 1024   # 25 MB — matches openvoiceui upload limit
MAX_PREVIEW_CHARS = 6000

ALLOWED_EXTRACT_EXTENSIONS = {
    '.pdf', '.docx', '.xlsx', '.pptx',
}

ALLOWED_CSV_EXTENSIONS = {'.csv', '.tsv'}

# Control characters: strip everything except \t \n \r
_CTRL_RE = re.compile(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]')


# ---------------------------------------------------------------------------
# Sanitization
# ---------------------------------------------------------------------------

def sanitize(text: str) -> str:
    """Strip control chars, collapse blank lines, cap at MAX_PREVIEW_CHARS."""
    text = _CTRL_RE.sub('', text)
    text = re.sub(r'\n{4,}', '\n\n\n', text)
    return text[:MAX_PREVIEW_CHARS].strip()


# ---------------------------------------------------------------------------
# Extractors
# ---------------------------------------------------------------------------

def extract_pdf(path: Path) -> tuple[str, dict]:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages):
        try:
            pages.append(page.extract_text() or '')
        except Exception:
            pages.append(f'[Page {i + 1}: extraction failed]')
    text = sanitize('\n\n'.join(pages))
    return text, {'pages': len(reader.pages)}


def extract_docx(path: Path) -> tuple[str, dict]:
    from docx import Document
    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(' | '.join(cells))
    text = sanitize('\n'.join(parts))
    return text, {'paragraphs': len(doc.paragraphs), 'tables': len(doc.tables)}


def extract_xlsx(path: Path) -> tuple[str, dict]:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    sheet_count = 0
    try:
        for sheet in wb.worksheets:
            sheet_count += 1
            parts.append(f'[Sheet: {sheet.title}]')
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    parts.append('\t'.join(cells))
    finally:
        wb.close()
    text = sanitize('\n'.join(parts))
    return text, {'sheets': sheet_count}


def extract_pptx(path: Path) -> tuple[str, dict]:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
        if slide_texts:
            parts.append(f'[Slide {i}]')
            parts.extend(slide_texts)
    text = sanitize('\n'.join(parts))
    return text, {'slides': len(prs.slides)}


# ---------------------------------------------------------------------------
# Routes
# ---------------------------------------------------------------------------

@app.route('/health')
def health():
    return jsonify({'status': 'ok', 'service': 'skill-runner'})


@app.route('/extract', methods=['POST'])
def extract():
    """
    Extract text from a document file.

    Accepts: multipart/form-data with:
      file     — the document bytes
      filename — original filename (used to determine type; optional, falls back to file.filename)

    Returns:
      {text, type, chars, ...type-specific meta}
    """
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400

    f = request.files['file']
    filename = request.form.get('filename') or f.filename or ''
    ext = Path(filename).suffix.lower()

    if ext not in ALLOWED_EXTRACT_EXTENSIONS:
        return jsonify({'error': f'Unsupported extension for extraction: {ext}'}), 415

    # Size check before writing to disk
    f.stream.seek(0, 2)
    size = f.stream.tell()
    f.stream.seek(0)
    if size > MAX_FILE_BYTES:
        return jsonify({'error': 'File too large (25 MB max)'}), 413

    # Write to a secure temp file
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp_path = Path(tmp.name)
        f.save(tmp_path)

    try:
        if ext == '.pdf':
            text, meta = extract_pdf(tmp_path)
        elif ext == '.docx':
            text, meta = extract_docx(tmp_path)
        elif ext == '.xlsx':
            text, meta = extract_xlsx(tmp_path)
        elif ext == '.pptx':
            text, meta = extract_pptx(tmp_path)
        else:
            return jsonify({'error': f'No extractor for {ext}'}), 415

    except Exception as exc:
        logger.exception('Extraction failed for %s', filename)
        return jsonify({'error': f'Extraction failed: {exc}'}), 500
    finally:
        try:
            tmp_path.unlink()
        except Exception:
            pass

    logger.info('Extracted %s → %d chars (%s)', filename, len(text), ext)
    return jsonify({'text': text, 'type': ext.lstrip('.'), 'chars': len(text), **meta})


@app.route('/analyze-csv', methods=['POST'])
def analyze_csv():
    """
    Run pandas summary analysis on a CSV/TSV file.

    Accepts: multipart/form-data with:
      file     — the CSV bytes
      filename — original filename

    Returns:
      {summary, rows, columns, dtypes, missing}
    """
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400

    f = request.files['file']
    filename = request.form.get('filename') or f.filename or ''
    ext = Path(filename).suffix.lower()

    if ext not in ALLOWED_CSV_EXTENSIONS:
        return jsonify({'error': f'Unsupported extension: {ext}. Send .csv or .tsv'}), 415

    f.stream.seek(0, 2)
    size = f.stream.tell()
    f.stream.seek(0)
    if size > MAX_FILE_BYTES:
        return jsonify({'error': 'File too large (25 MB max)'}), 413

    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp_path = Path(tmp.name)
        f.save(tmp_path)

    try:
        import pandas as pd
        sep = '\t' if ext == '.tsv' else ','
        df = pd.read_csv(tmp_path, sep=sep)

        rows, cols = df.shape
        dtypes = {col: str(dtype) for col, dtype in df.dtypes.items()}
        missing = {col: int(df[col].isna().sum()) for col in df.columns if df[col].isna().any()}

        # Build a readable text summary
        desc = df.describe(include='all').to_string()
        summary = sanitize(
            f'Rows: {rows}, Columns: {cols}\n\nColumn types:\n'
            + '\n'.join(f'  {c}: {t}' for c, t in dtypes.items())
            + f'\n\nStats:\n{desc}'
        )

    except Exception as exc:
        logger.exception('CSV analysis failed for %s', filename)
        return jsonify({'error': f'Analysis failed: {exc}'}), 500
    finally:
        try:
            tmp_path.unlink()
        except Exception:
            pass

    logger.info('Analyzed CSV %s → %d rows × %d cols', filename, rows, cols)
    return jsonify({
        'summary': summary,
        'rows': rows,
        'columns': cols,
        'dtypes': dtypes,
        'missing': missing,
    })


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 8900))
    logger.info('skill-runner starting on port %d', port)
    app.run(host='0.0.0.0', port=port, threaded=True)
